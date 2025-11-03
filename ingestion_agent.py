"""
Ingestion Agent - Document Processing Module
============================================

This agent handles:
1. Reading multiple document formats (PDF, DOCX, PPTX, CSV, TXT, MD)
2. Extracting text content from documents
3. Chunking text into manageable segments
4. Sending processed chunks to Retrieval Agent via MCP


"""

import os
import PyPDF2
from docx import Document
from pptx import Presentation
import logging
from mcp import send_mcp_message

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def extract_text_from_pdf(file):
    """
    Extract text content from PDF files.
    
    Args:
        file: File object or path to PDF
        
    Returns:
        str: Extracted text content
    """
    try:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        logger.info(f"Successfully extracted text from PDF ({len(text)} characters)")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        return ""


def extract_text_from_docx(file):
    """
    Extract text content from Word documents.
    
    Args:
        file: File object or path to DOCX
        
    Returns:
        str: Extracted text content
    """
    try:
        doc = Document(file)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        logger.info(f"Successfully extracted text from DOCX ({len(text)} characters)")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {str(e)}")
        return ""


def extract_text_from_pptx(file):
    """
    Extract text content from PowerPoint presentations.
    
    Args:
        file: File object or path to PPTX
        
    Returns:
        str: Extracted text content
    """
    try:
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        logger.info(f"Successfully extracted text from PPTX ({len(text)} characters)")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {str(e)}")
        return ""


def extract_text_from_txt(file):
    """
    Extract text content from plain text files.
    
    Args:
        file: File object or path to TXT/MD
        
    Returns:
        str: Extracted text content
    """
    try:
        if isinstance(file, str):
            with open(file, 'r', encoding='utf-8') as f:
                text = f.read()
        else:
            text = file.read().decode('utf-8')
        logger.info(f"Successfully extracted text from TXT ({len(text)} characters)")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from TXT: {str(e)}")
        return ""


def extract_text_from_csv(file):
    """
    Extract text content from CSV files.
    
    Args:
        file: File object or path to CSV
        
    Returns:
        str: Extracted text content (formatted as rows)
    """
    try:
        import csv
        import io
        
        if isinstance(file, str):
            with open(file, 'r', encoding='utf-8') as f:
                content = f.read()
        else:
            content = file.read().decode('utf-8')
        
        csv_reader = csv.reader(io.StringIO(content))
        text = "\n".join([", ".join(row) for row in csv_reader])
        logger.info(f"Successfully extracted text from CSV ({len(text)} characters)")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from CSV: {str(e)}")
        return ""


def chunk_text(text, chunk_size=500, overlap=50):
    """
    Split text into overlapping chunks for better context preservation.
    
    Args:
        text (str): Input text to chunk
        chunk_size (int): Size of each chunk in characters
        overlap (int): Number of overlapping characters between chunks
        
    Returns:
        list: List of text chunks
    """
    if not text or len(text) == 0:
        logger.warning("Empty text provided for chunking")
        return []
    
    chunks = []
    start = 0
    text_length = len(text)
    
    while start < text_length:
        end = start + chunk_size
        chunk = text[start:end]
        
        # Try to break at sentence or word boundary
        if end < text_length:
            # Look for sentence end
            last_period = chunk.rfind('.')
            last_newline = chunk.rfind('\n')
            last_space = chunk.rfind(' ')
            
            # Use the best boundary found
            break_point = max(last_period, last_newline, last_space)
            if break_point > chunk_size * 0.7:  # Only if boundary is reasonably close
                chunk = chunk[:break_point + 1]
                end = start + break_point + 1
        
        chunks.append(chunk.strip())
        start = end - overlap  # Create overlap
    
    logger.info(f"Text chunked into {len(chunks)} segments")
    return chunks


def process_document(file, filename):
    """
    Main processing function - coordinates document ingestion pipeline.
    
    Pipeline Steps:
    1. Identify file type
    2. Extract text using appropriate extractor
    3. Chunk text into segments
    4. Send to Retrieval Agent via MCP
    
    Args:
        file: File object to process
        filename (str): Name of the file
        
    Returns:
        bool: True if successful, False otherwise
    """
    logger.info(f"Starting document processing: {filename}")
    
    # Determine file type and extract text
    file_extension = filename.lower().split('.')[-1]
    
    text = ""
    if file_extension == "pdf":
        text = extract_text_from_pdf(file)
    elif file_extension == "docx":
        text = extract_text_from_docx(file)
    elif file_extension == "pptx":
        text = extract_text_from_pptx(file)
    elif file_extension in ["txt", "md"]:
        text = extract_text_from_txt(file)
    elif file_extension == "csv":
        text = extract_text_from_csv(file)
    else:
        logger.error(f"Unsupported file type: {file_extension}")
        return False
    
    if not text:
        logger.error(f"No text extracted from {filename}")
        return False
    
    # Chunk the text
    chunks = chunk_text(text, chunk_size=500, overlap=50)
    
    if not chunks:
        logger.error(f"No chunks created from {filename}")
        return False
    
    # Send to Retrieval Agent via MCP
    logger.info(f"Sending {len(chunks)} chunks to Retrieval Agent via MCP")
    message = {
        "sender": "IngestionAgent",
        "receiver": "RetrievalAgent",
        "payload": {
            "chunks": chunks,
            "document_name": filename,
            "metadata": {
                "file_type": file_extension,
                "chunk_count": len(chunks),
                "total_chars": len(text)
            }
        }
    }
    
    send_mcp_message(message)
    logger.info(f"Successfully processed document: {filename}")
    return True


if __name__ == "__main__":
    # Test the ingestion agent
    print("=" * 60)
    print("INGESTION AGENT - Document Processing Module")
    print("=" * 60)
    print("\nSupported formats: PDF, DOCX, PPTX, CSV, TXT, MD")
    print("\nThis module extracts text and chunks documents for RAG pipeline.")
    print("\nKey Features:")
    print("  ✓ Multi-format support")
    print("  ✓ Intelligent text chunking with overlap")
    print("  ✓ Sentence-boundary aware splitting")
    print("  ✓ MCP communication for agent coordination")
    print("=" * 60)
